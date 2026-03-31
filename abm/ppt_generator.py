"""
ppt_generator.py — Genera il report PPT campagna OOH.
Stile mediaitalia: sfondo nero/bianco, barra rossa verticale, footer nero.
"""
from __future__ import annotations
import os
import tempfile
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette mediaitalia ────────────────────────────────────────────────────────
RED        = RGBColor(0xCC, 0x00, 0x00)
BLACK      = RGBColor(0x00, 0x00, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
FOOTER_BG  = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY   = RGBColor(0x77, 0x77, 0x77)
ROW_ALT    = RGBColor(0xF8, 0xF8, 0xF8)
BORDER     = RGBColor(0xE0, 0xE0, 0xE0)

# ── Dimensioni slide 16:9 ─────────────────────────────────────────────────────
SLIDE_W  = Inches(13.33)
SLIDE_H  = Inches(7.5)
FOOTER_H = Inches(0.45)
MARGIN_L = Inches(0.65)
MARGIN_T = Inches(0.45)
BAR_W    = Inches(0.07)   # larghezza barra rossa
BAR_GAP  = Inches(0.18)   # gap tra barra e testo


# ── Helpers ───────────────────────────────────────────────────────────────────
def _rect(slide, left, top, width, height, fill: RGBColor, border: RGBColor = None):
    from pptx.util import Emu
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(0.5)
    else:
        shp.line.fill.background()
    return shp


def _textbox(slide, left, top, width, height, text, size, bold=False,
             color=BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = "Calibri"
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def _footer(slide, page_num: int, section: str = ""):
    top = SLIDE_H - FOOTER_H
    _rect(slide, 0, top, SLIDE_W, FOOTER_H, FOOTER_BG)
    _textbox(slide, Inches(0.3), top + Inches(0.07), Inches(3), FOOTER_H,
             "▶  mediaitalia", 10, bold=True, color=WHITE)
    if section:
        _textbox(slide, SLIDE_W / 2 - Inches(2), top + Inches(0.07), Inches(4), FOOTER_H,
                 section, 9, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)
    _textbox(slide, SLIDE_W - Inches(0.9), top + Inches(0.07), Inches(0.7), FOOTER_H,
             str(page_num), 10, color=WHITE, align=PP_ALIGN.RIGHT)


def _title(slide, text: str, top=None, dark=False, size=26):
    """Titolo con barra rossa verticale."""
    if top is None:
        top = MARGIN_T
    txt_color = WHITE if dark else BLACK
    bar_h = Inches(0.48)
    _rect(slide, MARGIN_L, top, BAR_W, bar_h, RED)
    _textbox(slide, MARGIN_L + BAR_W + BAR_GAP, top - Inches(0.03),
             SLIDE_W - MARGIN_L - BAR_W - BAR_GAP - Inches(0.4),
             bar_h + Inches(0.1), text, size, bold=True, color=txt_color, wrap=False)


# ── Slide 1: Cover ────────────────────────────────────────────────────────────
def _slide_cover(prs, cliente: str, nome_campagna: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, BLACK)

    # Barra rossa verticale alta
    _rect(slide, MARGIN_L, Inches(2.0), BAR_W, Inches(1.5), RED)

    # Titolo principale
    _textbox(slide, MARGIN_L + BAR_W + BAR_GAP, Inches(2.0),
             Inches(9), Inches(0.8),
             "REPORT CAMPAGNA OOH", 38, bold=True, color=WHITE)

    # Cliente
    _textbox(slide, MARGIN_L + BAR_W + BAR_GAP, Inches(2.9),
             Inches(9), Inches(0.6),
             cliente.upper(), 22, bold=True, color=WHITE)

    # Nome campagna
    _textbox(slide, MARGIN_L + BAR_W + BAR_GAP, Inches(3.55),
             Inches(9), Inches(0.45),
             nome_campagna, 14, color=RGBColor(0xBB, 0xBB, 0xBB))

    # Data
    mese_anno = datetime.now().strftime("%B %Y").capitalize()
    _textbox(slide, MARGIN_L + BAR_W + BAR_GAP, Inches(4.1),
             Inches(4), Inches(0.35),
             mese_anno, 11, color=RGBColor(0x88, 0x88, 0x88))

    _footer(slide, 1)


# ── Slide 2: Sommario campagna ────────────────────────────────────────────────
def _slide_sommario(prs, info: dict, metriche_agg: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    _title(slide, "Piano Campagna")

    # ── KPI boxes (2 righe × 3 colonne) ──────────────────────────────────────
    kpis = [
        ("Comuni",       f"{info['n_comuni']}"),
        ("Impianti",     f"{info['n_impianti']}"),
        ("OTS Totali",   f"{metriche_agg['ots_totali']:,}"),
        ("Reach",        f"{metriche_agg['reach_univoco']:,}"),
        ("GRP Medio",    f"{metriche_agg['grp_medio']:.1f}"),
        ("Copertura",    f"{metriche_agg['copertura_media']:.1f}%"),
    ]
    bw, bh = Inches(2.0), Inches(1.3)
    gap    = Inches(0.22)
    t0     = Inches(1.3)

    for i, (label, value) in enumerate(kpis):
        col = i % 3
        row = i // 3
        lft = MARGIN_L + col * (bw + gap)
        top = t0 + row * (bh + gap)
        _rect(slide, lft, top, bw, bh, LIGHT_GRAY, BORDER)
        _textbox(slide, lft + Inches(0.15), top + Inches(0.12),
                 bw - Inches(0.2), Inches(0.72), value, 26, bold=True, color=BLACK)
        _textbox(slide, lft + Inches(0.15), top + Inches(0.88),
                 bw - Inches(0.2), Inches(0.3), label, 9, color=MID_GRAY)

    # ── Composizione piano (colonna destra) ───────────────────────────────────
    rx = MARGIN_L + Inches(7.2)
    _textbox(slide, rx, Inches(1.3), Inches(5), Inches(0.4),
             "Composizione Piano", 13, bold=True, color=BLACK)

    for i, (tipo, count) in enumerate(info.get("impianti_per_tipo", {}).items()):
        ty = Inches(1.85) + i * Inches(0.48)
        _rect(slide, rx, ty + Inches(0.08), Inches(0.06), Inches(0.28), RED)
        _textbox(slide, rx + Inches(0.18), ty, Inches(4.5), Inches(0.4),
                 f"{tipo.capitalize()}:  {count}", 12, color=BLACK)

    _footer(slide, 2, "Sommario")


# ── Slide 3+: Tabella metriche ────────────────────────────────────────────────
def _slide_tabella(prs, impianti_dict: dict, metriche: dict, start_page: int):
    items    = sorted(metriche.items(),
                      key=lambda x: x[1]["ots_totali"], reverse=True)
    MAX_ROWS = 17

    cols   = ["ID Impianto", "Tipo",  "Indirizzo",     "OTS",   "Reach",   "Freq.", "GRP",  "Cop.%"]
    widths = [Inches(1.1),   Inches(0.95), Inches(3.6), Inches(0.95), Inches(1.1), Inches(0.75), Inches(0.8), Inches(0.75)]
    aligns = [PP_ALIGN.LEFT] * 3 + [PP_ALIGN.RIGHT] * 5
    ROW_H  = Inches(0.3)
    HDR_H  = Inches(0.33)

    for slide_idx, start in enumerate(range(0, len(items), MAX_ROWS)):
        batch = items[start:start + MAX_ROWS]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)

        suffix = f" ({slide_idx + 1})" if len(items) > MAX_ROWS else ""
        _title(slide, f"Metriche per Impianto{suffix}")

        # posizioni colonne
        lefts = []
        x = MARGIN_L
        for w in widths:
            lefts.append(x)
            x += w

        hdr_top = Inches(1.1)

        # Header
        for col, left, width, aln in zip(cols, lefts, widths, aligns):
            _rect(slide, left, hdr_top, width, HDR_H, BLACK)
            _textbox(slide, left + Inches(0.05), hdr_top + Inches(0.05),
                     width - Inches(0.07), HDR_H, col, 9,
                     bold=True, color=WHITE, align=aln, wrap=False)

        # Righe
        for ri, (imp_id, m) in enumerate(batch):
            imp     = impianti_dict.get(imp_id)
            row_top = hdr_top + HDR_H + ri * ROW_H
            fill    = ROW_ALT if ri % 2 == 0 else WHITE
            ind     = (imp.indirizzo[:38] + "…") if imp and len(imp.indirizzo) > 38 else (imp.indirizzo if imp else "")

            values = [
                imp_id,
                imp.tipo.capitalize() if imp else "",
                ind,
                f"{m['ots_totali']:,}",
                f"{m['reach_univoco']:,}",
                f"{m['frequency_media']:.1f}",
                f"{m['grp']:.1f}",
                f"{m['copertura_pct']:.2f}%",
            ]
            for val, left, width, aln in zip(values, lefts, widths, aligns):
                _rect(slide, left, row_top, width, ROW_H, fill, BORDER)
                _textbox(slide, left + Inches(0.05), row_top + Inches(0.04),
                         width - Inches(0.07), ROW_H, str(val), 9,
                         color=BLACK, align=aln, wrap=False)

        _footer(slide, start_page + slide_idx, "Metriche Impianti")


# ── Entry point ───────────────────────────────────────────────────────────────
def genera_ppt(data: dict, output_path: str = None) -> str:
    """
    Genera il PPT report OOH.

    Args:
        data: {
            'campagna':  dict da load_campaign (con cliente, nome_campagna, ecc.)
            'risultati': dict impianto_id → metriche
            'impianti':  dict impianto_id → ImpiantoABM
        }
        output_path: path di salvataggio (default: temp file)

    Returns:
        path del file .pptx generato
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    camp      = data.get("campagna", {})
    risultati = data.get("risultati", {})
    impianti  = data.get("impianti", {})

    cliente       = camp.get("cliente", "")
    nome_campagna = camp.get("nome_campagna", f"Campagna {camp.get('id_campagna', '')}")

    # Metriche aggregate
    if risultati:
        ots_tot    = sum(m["ots_totali"] for m in risultati.values())
        reach_max  = max(m["reach_univoco"] for m in risultati.values())
        grp_medio  = sum(m["grp"] for m in risultati.values()) / len(risultati)
        cop_media  = sum(m["copertura_pct"] for m in risultati.values()) / len(risultati)
    else:
        ots_tot = reach_max = grp_medio = cop_media = 0

    impianti_per_tipo: dict[str, int] = {}
    for imp in impianti.values():
        impianti_per_tipo[imp.tipo] = impianti_per_tipo.get(imp.tipo, 0) + 1

    info = {
        "n_comuni":          len(camp.get("comuni_istat", [])),
        "n_impianti":        len(impianti),
        "impianti_per_tipo": impianti_per_tipo,
    }
    metriche_agg = {
        "ots_totali":      ots_tot,
        "reach_univoco":   reach_max,
        "grp_medio":       grp_medio,
        "copertura_media": cop_media,
    }

    _slide_cover(prs, cliente, nome_campagna)
    _slide_sommario(prs, info, metriche_agg)
    _slide_tabella(prs, impianti, risultati, start_page=3)

    if output_path is None:
        output_path = os.path.join(
            tempfile.gettempdir(),
            f"report_ooh_{camp.get('id_campagna', 'x')}.pptx"
        )
    prs.save(output_path)
    return output_path
